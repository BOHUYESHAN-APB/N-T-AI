"""
PPT Generator Tool
Converts HTML content (with Tailwind-like classes) into a PowerPoint (.pptx) presentation.
Inspired by free-OKC's slides_generator tool (MIT License).
"""

import os
import datetime
from pathlib import Path
from typing import List, Tuple, Dict, Optional
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PPTGenerator:
    def __init__(self, output_dir: str = "generated_docs/ppt"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def _parse_slides(self, html: str) -> List[BeautifulSoup]:
        """Parses HTML and extracts slide elements."""
        soup = BeautifulSoup(html, "html.parser")
        # Assuming slides are wrapped in a container with class 'ppt-slide'
        slides = soup.select(".ppt-slide")
        if not slides:
            # Fallback: treat top-level sections or divs as slides if no specific class found
            slides = soup.find_all("section")
            if not slides:
                 # If still nothing, wrap the whole body as one slide (emergency fallback)
                 return [soup]
        return slides

    def _extract_slide_content(self, slide_markup: BeautifulSoup, index: int) -> Dict:
        """Extracts title, paragraphs, and list items from a slide's HTML."""
        
        # Extract Title
        title_tag = slide_markup.find(["h1", "h2", "h3"])
        title = title_tag.get_text(strip=True) if title_tag else f"Slide {index + 1}"

        # Extract Content (Paragraphs and Lists)
        # We iterate through children to preserve order
        content_elements = []
        
        for element in slide_markup.find_all(["p", "li"]):
            text = element.get_text(strip=True)
            if not text:
                continue
            
            tag_name = element.name
            if tag_name == "li":
                content_elements.append({"type": "bullet", "text": text})
            else:
                content_elements.append({"type": "text", "text": text})

        return {"title": title, "elements": content_elements}

    def _add_textbox(self, slide, text: str, left: float, top: float, width: float, height: float, font_size: int = 18, is_bullet: bool = False):
        """Adds a textbox to the slide."""
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        
        if is_bullet:
            p.level = 0 # First level bullet

    def generate_ppt(self, html_content: str, filename: str = None) -> str:
        """
        Generates a PPTX file from the provided HTML content.
        
        Args:
            html_content: The HTML string describing the slides.
            filename: Optional filename. If not provided, a timestamped one will be generated.
            
        Returns:
            The absolute path to the generated PPTX file.
        """
        if not filename:
            timestamp = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
            filename = f"presentation_{timestamp}.pptx"
            
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        output_path = self.output_dir / filename
        
        presentation = Presentation()
        # Use a blank layout (usually index 6 in default template)
        blank_layout_index = 6 
        
        slides_markup = self._parse_slides(html_content)
        
        for i, slide_markup in enumerate(slides_markup):
            # Check if layout exists, otherwise use 0
            if len(presentation.slide_layouts) > blank_layout_index:
                slide_layout = presentation.slide_layouts[blank_layout_index]
            else:
                slide_layout = presentation.slide_layouts[0]
                
            slide = presentation.slides.add_slide(slide_layout)
            
            content = self._extract_slide_content(slide_markup, i)
            
            # 1. Add Title
            # Centered title at the top
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0))
            title_tf = title_shape.text_frame
            title_p = title_tf.paragraphs[0]
            title_p.text = content["title"]
            title_p.font.size = Pt(32)
            title_p.font.bold = True
            title_p.alignment = PP_ALIGN.CENTER

            # 2. Add Content
            # Start position for content
            current_top = 1.5
            line_height = 0.5
            
            for element in content["elements"]:
                text = element["text"]
                is_bullet = element["type"] == "bullet"
                
                # Simple layout logic: just stack them vertically
                # For a more advanced version, we could calculate text height
                
                if is_bullet:
                    text = f"• {text}"
                
                self._add_textbox(
                    slide, 
                    text, 
                    left=1.0, 
                    top=current_top, 
                    width=8.0, 
                    height=line_height, 
                    font_size=20 if is_bullet else 24,
                    is_bullet=is_bullet
                )
                
                current_top += 0.6 # Increment vertical position
                
                # Simple page break logic (if too much content)
                if current_top > 7.0:
                    # Create a new slide if overflow (simplified)
                     slide = presentation.slides.add_slide(slide_layout)
                     current_top = 1.5
        
        presentation.save(str(output_path))
        return str(output_path.absolute())

if __name__ == "__main__":
    # Simple test
    generator = PPTGenerator()
    sample_html = """
    <div class="ppt-slide">
        <h1>Welcome to N-T-AI</h1>
        <p>This is a generated presentation.</p>
        <li>Feature 1: AI Chat</li>
        <li>Feature 2: PPT Generation</li>
    </div>
    <div class="ppt-slide">
        <h1>Architecture</h1>
        <p>Backend: FastAPI</p>
        <p>Frontend: Flutter</p>
    </div>
    """
    print(generator.generate_ppt(sample_html))
