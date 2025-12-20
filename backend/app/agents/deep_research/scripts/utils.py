import re
import json
import html
from typing import List, Dict, Any, Optional
from pathlib import Path
from app.skills.common.document_skill.scripts.office_processor import OfficeProcessor

def extract_json_from_text(text: str) -> Optional[Dict[str, Any]]:
    """
    Extracts JSON object from text, handling markdown code blocks.
    """
    text = (text or "").strip()
    # Remove markdown code blocks
    match = re.search(r"```(?:json)?\s*([\s\S]*?)\s*```", text)
    if match:
        text = match.group(1)
    
    try:
        # Find the first '{' and last '}'
        start = text.find('{')
        end = text.rfind('}')
        if start != -1 and end != -1:
            text = text[start:end+1]
        return json.loads(text)
    except json.JSONDecodeError:
        return None

def extract_file_content(context_files: List[Dict[str, Any]], office_processor: Any) -> str:
    revision_context_text = ""
    extracted_parts: List[str] = []
    
    if not context_files:
        return ""

    for item in context_files:
        if not isinstance(item, dict):
            continue
        path = str(item.get("path") or "").strip()
        title = str(item.get("title") or "").strip()
        if not path:
            continue
        try:
            p = Path(path)
            if not p.exists():
                p = Path("app/static/reports") / p.name
            if not p.exists():
                continue

            suffix = p.suffix.lower()
            extracted = ""
            if suffix == ".pdf":
                extracted = office_processor.extract_text_from_pdf(str(p))
            elif suffix == ".docx":
                from docx import Document as DocxDocument
                doc = DocxDocument(str(p))
                extracted = "\n".join([para.text for para in doc.paragraphs if para.text])
            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(p))
                lines: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if getattr(shape, "has_text_frame", False) and shape.text_frame:
                            t = (shape.text_frame.text or "").strip()
                            if t:
                                lines.append(t)
                extracted = "\n".join(lines)
            elif suffix in (".xlsx", ".xls"):
                try:
                    import pandas as pd
                    df = pd.read_excel(str(p))
                    extracted = df.head(50).to_csv(index=False)
                except Exception:
                    extracted = ""

            extracted = (extracted or "").strip()
            if extracted:
                extracted = extracted[:8000]
                extracted_parts.append(f"=== {title or p.name} ({p.name}) ===\n{extracted}")
        except Exception:
            continue

    if extracted_parts:
        revision_context_text = "\n\n".join(extracted_parts)
    
    return revision_context_text

def normalize_user_input(user_input: str) -> str:
    text = (user_input or "").strip()
    if not text:
        return ""
    text = re.sub(r"^User Request:\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"^\s*Topic\s*&\s*Constraints.*?:\s*", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.replace("原始需求:", "")
    text = text.replace("补充信息:", "")
    text = text.replace("未填写项请AI自行补全，并继续执行。", "")
    text = text.replace("用户60秒未响应，请基于现有信息自行补全缺失项并继续执行。", "")
    text = text.replace("用户选择跳过补充信息，请基于现有信息自行补全缺失项并继续执行。", "")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text

def truncate_text(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + "\n...[truncated]..."

def ensure_html(content: str, office_processor: Any) -> str:
    """
    Ensures the content is HTML. If it's JSON or Markdown, converts it.
    Requires an office_processor instance for markdown conversion.
    """
    text = (content or "").strip()
    if not text:
        return ""

    # Try to parse as JSON first if it looks like JSON
    if text.startswith("{") or text.startswith("["):
        try:
            data = json.loads(text)
            # If valid JSON, try to extract meaningful content
            parts = []
            if isinstance(data, list):
                for item in data:
                    parts.append(str(item))
            elif isinstance(data, dict):
                # Check for common document structure fields
                if "title" in data: parts.append(f"# {data['title']}")
                
                content_body = data.get("content") or data.get("body") or data.get("text")
                sections = data.get("sections")
                
                if sections and isinstance(sections, list):
                    for sec in sections:
                        if isinstance(sec, dict):
                            if "title" in sec: parts.append(f"## {sec['title']}")
                            if "content" in sec: parts.append(str(sec['content']))
                        else:
                            parts.append(str(sec))
                elif content_body:
                    parts.append(str(content_body))
                else:
                    # No obvious content fields, just dump keys/values
                    for k, v in data.items():
                        if k not in ["title", "format"]:
                            parts.append(f"**{k}**: {v}")
            
            if parts:
                text = "\n\n".join(parts)
        except:
            pass # Not valid JSON, treat as text

    if re.search(r"</?(h\d|p|ul|ol|li|div|table|img)\b", text, flags=re.IGNORECASE):
        return text
    return office_processor.markdown_to_html(text)

def wrap_preview_html(kind: str, html_content: str) -> str:
    body = (html_content or "").strip()
    if not body:
        body = "<p>(空)</p>"
    base_css = (
        "body{font-family:Arial,'Microsoft YaHei',sans-serif;line-height:1.6;padding:24px;"
        "background:#f7f9fc;color:#0f172a}"
        "h1,h2,h3{color:#0b3a6a}"
        "table{border-collapse:collapse;width:100%}"
        "th,td{border:1px solid #e2e8f0;padding:8px}"
        "code,pre{font-family:Consolas,monospace}"
    )
    if kind == "ppt":
        base_css += (
            ".ppt-slide{background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;"
            "padding:24px;margin:18px auto;max-width:980px;box-shadow:0 6px 24px rgba(15,23,42,.08)}"
            ".ppt-slide h1{font-size:32px;text-align:center;margin:0 0 12px}"
            ".ppt-slide h2{font-size:22px;margin:0 0 10px}"
            ".ppt-slide ul{padding-left:22px}"
            ".ppt-slide table.chart{margin-top:10px}"
        )
    return f"<html><head><meta charset='utf-8'><style>{base_css}</style></head><body>{body}</body></html>"

def detect_requested_formats(user_input: str) -> List[str]:
    text = (user_input or "").lower()
    requested: List[str] = []
    if any(k in text for k in ["docx", "word", "文档", "报告", ".doc"]):
        requested.append("docx")
    if any(k in text for k in ["pptx", "ppt", "powerpoint", "演示", "幻灯片"]):
        requested.append("ppt")
    if any(k in text for k in ["pdf"]):
        requested.append("pdf")
    if any(k in text for k in ["xlsx", "excel", ".xlsx", "生成excel", "导出excel", "excel表", "电子表格"]):
        requested.append("excel")
    if not requested:
        requested.append("docx")
    deduped: List[str] = []
    for f in requested:
        if f not in deduped:
            deduped.append(f)
    return deduped

def normalize_output_format(raw_format: Any) -> str:
    output_format = str(raw_format or "docx").strip().lower()
    if output_format in ["pptx", "powerpoint"]:
        return "ppt"
    if output_format in ["doc", "word"]:
        return "docx"
    if output_format in ["xlsx"]:
        return "excel"
    return output_format

def parse_writer_output_files(writer_response: str, user_input: str) -> List[Dict[str, Any]]:
    output_data_list: List[Dict[str, Any]] = []
    try:
        clean_json = writer_response
        if "```json" in clean_json:
            clean_json = clean_json.split("```json")[1].split("```")[0].strip()
        elif "```" in clean_json:
            clean_json = clean_json.split("```")[1].split("```")[0].strip()
        parsed = json.loads(clean_json)

        if isinstance(parsed, dict) and "files" in parsed and isinstance(parsed["files"], list):
            output_data_list = parsed["files"]
        elif isinstance(parsed, dict) and "format" in parsed:
            output_data_list = [parsed]
    except Exception:
        output_data_list = []

    if not output_data_list:
        filename_prefix = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", (user_input or "Report")[:20]).strip("_")
        output_data_list = [{
            "format": "docx",
            "content": writer_response,
            "filename": f"Report_{filename_prefix}.docx"
        }]

    return output_data_list
