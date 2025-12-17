"""
Deep Research Agent
A hierarchical agent that plans tasks and uses specialized tools (Sandbox, Office Suite).
Inspired by OpenManus (ReAct + Planning) and DeepResearchAgent (Hierarchy).
"""

from typing import List, Dict, Any, Optional, Tuple
from pydantic import BaseModel

from app.services.sandbox_service import sandbox_service
from app.tools.office_suite import PPTGenerator, WordGenerator, PDFGenerator, OfficeProcessor
from app.tools.office_suite.excel_generator import excel_generator
from app.services.search_service import SearchService
from app.services.vision_service import vision_service
from app.tools.academic_search import academic_search
from app.services.prompt_manager import prompt_manager

class TaskPlan(BaseModel):
    steps: List[str]
    current_step_index: int = 0

import datetime
import html
import os
import re
from pathlib import Path
from openai import AsyncOpenAI
import json
import httpx

class DeepResearchAgent:
    def __init__(self, model_config: Dict[str, Any]):
        """
        Args:
            model_config: Config for the LLM (API Key, Base URL, Model Name).
                          Structure: {"planner": {...}, "researcher": {...}, "writer": {...}}
        """
        self.config = model_config
        self.sandbox_session_id = sandbox_service.create_session()
        
        output_dir = "app/static/reports"
        self.ppt_generator = PPTGenerator(output_dir=output_dir)
        self.word_generator = WordGenerator(output_dir=output_dir)
        self.pdf_generator = PDFGenerator(output_dir=output_dir)
        self.office_processor = OfficeProcessor(output_dir=output_dir)
        self.search_service = SearchService()
        
        # Memory specific to this agent instance
        self.memory: List[Dict[str, str]] = []
        self.plan: Optional[TaskPlan] = None

    def _normalize_user_input(self, user_input: str) -> str:
        text = (user_input or "").strip()
        if not text:
            return ""
        text = re.sub(r"^User Request:\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"^\s*Topic\s*&\s*Constraints.*?:\s*", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\n{3,}", "\n\n", text)
        text = text.replace("原始需求:", "")
        text = text.replace("补充信息:", "")
        text = text.replace("未填写项请AI自行补全，并继续执行。", "")
        text = text.replace("用户60秒未响应，请基于现有信息自行补全缺失项并继续执行。", "")
        text = text.replace("用户选择跳过补充信息，请基于现有信息自行补全缺失项并继续执行。", "")
        text = re.sub(r"\n{3,}", "\n\n", text).strip()
        return text

    def _get_sandbox_workspace_dir(self) -> Path:
        session = sandbox_service.get_session(self.sandbox_session_id)
        if session and getattr(session, "workspace_dir", None):
            return Path(session.workspace_dir)
        return Path("workspace") / self.sandbox_session_id

    def _write_workspace_text(self, relative_path: str, text: str) -> str:
        workspace_dir = self._get_sandbox_workspace_dir()
        target_path = workspace_dir / relative_path
        target_path.parent.mkdir(parents=True, exist_ok=True)
        target_path.write_text(text, encoding="utf-8")
        return str(target_path)

    def _truncate_text(self, text: str, limit: int) -> str:
        if len(text) <= limit:
            return text
        return text[:limit] + "\n...[truncated]..."

    def _ensure_html(self, content: str) -> str:
        text = (content or "").strip()
        if not text:
            return ""

        # Try to parse as JSON first if it looks like JSON
        if text.startswith("{") or text.startswith("["):
            try:
                data = json.loads(text)
                # If valid JSON, try to extract meaningful content
                parts = []
                if isinstance(data, list):
                    for item in data:
                        parts.append(str(item))
                elif isinstance(data, dict):
                    # Check for common document structure fields
                    if "title" in data: parts.append(f"# {data['title']}")
                    
                    content_body = data.get("content") or data.get("body") or data.get("text")
                    sections = data.get("sections")
                    
                    if sections and isinstance(sections, list):
                        for sec in sections:
                            if isinstance(sec, dict):
                                if "title" in sec: parts.append(f"## {sec['title']}")
                                if "content" in sec: parts.append(str(sec['content']))
                            else:
                                parts.append(str(sec))
                    elif content_body:
                        parts.append(str(content_body))
                    else:
                        # No obvious content fields, just dump keys/values
                        for k, v in data.items():
                            if k not in ["title", "format"]:
                                parts.append(f"**{k}**: {v}")
                
                if parts:
                    text = "\n\n".join(parts)
            except:
                pass # Not valid JSON, treat as text

        if re.search(r"</?(h\d|p|ul|ol|li|div|table|img)\b", text, flags=re.IGNORECASE):
            return text
        return self.office_processor.markdown_to_html(text)

    def _wrap_preview_html(self, kind: str, html_content: str) -> str:
        body = (html_content or "").strip()
        if not body:
            body = "<p>(空)</p>"
        base_css = (
            "body{font-family:Arial,'Microsoft YaHei',sans-serif;line-height:1.6;padding:24px;"
            "background:#f7f9fc;color:#0f172a}"
            "h1,h2,h3{color:#0b3a6a}"
            "table{border-collapse:collapse;width:100%}"
            "th,td{border:1px solid #e2e8f0;padding:8px}"
            "code,pre{font-family:Consolas,monospace}"
        )
        if kind == "ppt":
            base_css += (
                ".ppt-slide{background:#ffffff;border:1px solid #e2e8f0;border-radius:12px;"
                "padding:24px;margin:18px auto;max-width:980px;box-shadow:0 6px 24px rgba(15,23,42,.08)}"
                ".ppt-slide h1{font-size:32px;text-align:center;margin:0 0 12px}"
                ".ppt-slide h2{font-size:22px;margin:0 0 10px}"
                ".ppt-slide ul{padding-left:22px}"
                ".ppt-slide table.chart{margin-top:10px}"
            )
        return f"<html><head><meta charset='utf-8'><style>{base_css}</style></head><body>{body}</body></html>"

    async def _skywork_generate_file_url(self, file_type: str, query: str) -> Optional[str]:
        api_key = os.getenv("SKYWORK_API_KEY")
        if not api_key:
            return None

        gen_url = "https://api-cn.tiangong.cn/infra/tool/generate_file"
        params = {"api_key": api_key, "query": query, "file_type": file_type}
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.post(gen_url, params=params)
            resp.raise_for_status()
            payload = resp.json()
            if payload.get("code") != 200:
                return None
            data = payload.get("data") or {}
            url = (data.get("url") or "").strip()
            return url or None

    async def _download_to_path(self, url: str, target_path: str) -> str:
        os.makedirs(os.path.dirname(target_path), exist_ok=True)
        async with httpx.AsyncClient(timeout=180) as client:
            resp = await client.get(url, follow_redirects=True)
            resp.raise_for_status()
            Path(target_path).write_bytes(resp.content)
        return target_path

    def _build_fallback_ppt_html(self, user_input: str, research_findings: str) -> str:
        title = (user_input or "Research Report").strip()
        title = title[:60] if title else "Research Report"
        brief = self._truncate_text(research_findings.strip(), 1200)
        brief = html.escape(brief)
        safe_title = html.escape(title)

        slides = [
            f"<div class=\"ppt-slide\"><h1>{safe_title}</h1><p>图文并茂 / 含表格与柱状图</p></div>",
            f"<div class=\"ppt-slide\"><h2>关键要点</h2><p>{brief}</p></div>",
            (
                "<div class=\"ppt-slide\">"
                "<h2>数据概览（示例）</h2>"
                "<p>如模型未返回统计数据，此处为占位示例，建议在下一轮检索补齐。</p>"
                "<table class=\"chart\">"
                "<tr><th>类别</th><th>数量</th></tr>"
                "<tr><td>产品发布</td><td>6</td></tr>"
                "<tr><td>融资/财报</td><td>4</td></tr>"
                "<tr><td>模型/论文</td><td>8</td></tr>"
                "<tr><td>政策/监管</td><td>3</td></tr>"
                "</table>"
                "</div>"
            ),
            (
                "<div class=\"ppt-slide\">"
                "<h2>来源与证据</h2>"
                "<p>详见沙箱临时笔记文件（notes/step_*.md）。</p>"
                "</div>"
            ),
        ]
        return "\n".join(slides)

    def _build_search_queries(self, base_query: str, user_input: str) -> List[str]:
        base_query = (base_query or "").strip()
        if not base_query:
            base_query = user_input.strip()

        candidates = [
            base_query,
            f"{base_query} 最新",
            f"{base_query} 数据 统计",
            f"{base_query} 报告 2025",
            f"{base_query} 产品 发布",
        ]
        fallback_candidates = [
            f"{base_query} 论文 arxiv",
            f"{base_query} site:arxiv.org",
            f"{base_query} 知网 cnki",
            f"{base_query} SCI Web of Science",
        ]

        seen = set()
        queries: List[str] = []
        for q in candidates:
            q = re.sub(r"\s+", " ", (q or "").strip())
            if not q:
                continue
            if q in seen:
                continue
            seen.add(q)
            queries.append(q)

        for q in fallback_candidates:
            if len(queries) >= 4:
                break
            q = re.sub(r"\s+", " ", (q or "").strip())
            if not q or q in seen:
                continue
            seen.add(q)
            queries.append(q)

        if len(queries) < 3:
            for suffix in ["趋势", "对比", "行业 研究"]:
                if len(queries) >= 3:
                    break
                q = re.sub(r"\s+", " ", f"{base_query} {suffix}".strip())
                if q and q not in seen:
                    seen.add(q)
                    queries.append(q)

        return queries[: max(3, min(4, len(queries)))]

    async def _can_access_url(self, url: str) -> bool:
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
            }
            async with httpx.AsyncClient(timeout=6.0, follow_redirects=True, verify=False) as client:
                resp = await client.get(url, headers=headers)
            return resp.status_code >= 200 and resp.status_code < 500
        except Exception:
            return False

    def _detect_requested_formats(self, user_input: str) -> List[str]:
        text = (user_input or "").lower()
        requested: List[str] = []
        if any(k in text for k in ["docx", "word", "文档", "报告", ".doc"]):
            requested.append("docx")
        if any(k in text for k in ["pptx", "ppt", "powerpoint", "演示", "幻灯片"]):
            requested.append("ppt")
        if any(k in text for k in ["pdf"]):
            requested.append("pdf")
        if any(k in text for k in ["xlsx", "excel", ".xlsx", "生成excel", "导出excel", "excel表", "电子表格"]):
            requested.append("excel")
        if not requested:
            requested.append("docx")
        deduped: List[str] = []
        for f in requested:
            if f not in deduped:
                deduped.append(f)
        return deduped

    def _normalize_output_format(self, raw_format: Any) -> str:
        output_format = str(raw_format or "docx").strip().lower()
        if output_format in ["pptx", "powerpoint"]:
            return "ppt"
        if output_format in ["doc", "word"]:
            return "docx"
        if output_format in ["xlsx"]:
            return "excel"
        return output_format

    def _parse_writer_output_files(self, writer_response: str, user_input: str) -> List[Dict[str, Any]]:
        output_data_list: List[Dict[str, Any]] = []
        try:
            clean_json = writer_response
            if "```json" in clean_json:
                clean_json = clean_json.split("```json")[1].split("```")[0].strip()
            elif "```" in clean_json:
                clean_json = clean_json.split("```")[1].split("```")[0].strip()
            parsed = json.loads(clean_json)

            if isinstance(parsed, dict) and "files" in parsed and isinstance(parsed["files"], list):
                output_data_list = parsed["files"]
            elif isinstance(parsed, dict) and "format" in parsed:
                output_data_list = [parsed]
        except Exception:
            output_data_list = []

        if not output_data_list:
            filename_prefix = re.sub(r"[^\w\u4e00-\u9fff\-]+", "_", (user_input or "Report")[:20]).strip("_")
            output_data_list = [{
                "format": "docx",
                "content": writer_response,
                "filename": f"Report_{filename_prefix}.docx"
            }]

        return output_data_list

    def _get_client(self, role: str) -> AsyncOpenAI:
        """Creates an OpenAI client for the specific role."""
        role_config = self.config.get(role) or {}
        
        # Fallback to env vars if not provided
        api_key = role_config.get("api_key") or os.getenv("OPENAI_API_KEY")
        base_url = role_config.get("base_url") or os.getenv("OPENAI_BASE_URL")
        
        if not api_key:
            # Last resort for demo: raise error or use placeholder
            # For robustness, we might try to grab from settings singleton if possible, 
            # but usually it's passed in.
            from app.core.config import settings
            api_key = settings.OPENAI_API_KEY
            base_url = settings.OPENAI_BASE_URL
        
        if not api_key:
             raise ValueError(f"No API Key provided for role: {role}")

        return AsyncOpenAI(api_key=api_key, base_url=base_url)

    def _get_model(self, role: str) -> str:
        role_config = self.config.get(role) or {}
        return role_config.get("model") or "gpt-3.5-turbo"

    async def _call_llm(self, role: str, messages: List[Dict[str, str]], stream: bool = False) -> str:
        client = self._get_client(role)
        model = self._get_model(role)
        
        response = await client.chat.completions.create(
            model=model,
            messages=messages,
            stream=False # For internal logic we might not need stream
        )
        return response.choices[0].message.content

    def _append_research_log(self, content: str):
        """Appends content to the main research log markdown file."""
        try:
            log_path = self._get_sandbox_workspace_dir() / "research_log.md"
            with open(log_path, "a", encoding="utf-8") as f:
                f.write(content + "\n\n")
        except Exception:
            pass

    async def _generate_project_title(self, user_input: str) -> str:
        """Generates a short, concise title for the project."""
        try:
            sys_prompt = "You are a helpful assistant. Generate a short, concise title (max 15 chars) for the user's research request. No quotes, no markdown."
            response = await self._call_llm("planner", [
                {"role": "system", "content": sys_prompt},
                {"role": "user", "content": f"Request: {user_input}"}
            ])
            return response.strip()
        except:
            return "New Project"

    async def run_stream(self, user_input: str, context_files: List[Dict] = []):
        """
        Stream version of run. Yields events.
        """
        current_date = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        user_input = self._normalize_user_input(user_input)
        
        # 0. Generate Title & Init Log
        project_title = await self._generate_project_title(user_input)
        yield {
            "type": "metadata",
            "session_id": self.sandbox_session_id,
            "title": project_title
        }
        
        self._append_research_log(f"# Research Log: {project_title}\n\n**Date**: {current_date}\n**Request**: {user_input}\n\n---\n")

        revision_context_text = ""
        revision_context_titles: List[str] = []
        if context_files:
            extracted_parts: List[str] = []
            for item in context_files:
                if not isinstance(item, dict):
                    continue
                path = str(item.get("path") or "").strip()
                title = str(item.get("title") or "").strip()
                if not path:
                    continue
                try:
                    p = Path(path)
                    if not p.exists():
                        p = Path("app/static/reports") / p.name
                    if not p.exists():
                        continue

                    suffix = p.suffix.lower()
                    extracted = ""
                    if suffix == ".pdf":
                        extracted = self.office_processor.extract_text_from_pdf(str(p))
                    elif suffix == ".docx":
                        from docx import Document as DocxDocument
                        doc = DocxDocument(str(p))
                        extracted = "\n".join([para.text for para in doc.paragraphs if para.text])
                    elif suffix == ".pptx":
                        from pptx import Presentation
                        prs = Presentation(str(p))
                        lines: List[str] = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                                    t = (shape.text_frame.text or "").strip()
                                    if t:
                                        lines.append(t)
                        extracted = "\n".join(lines)
                    elif suffix in (".xlsx", ".xls"):
                        try:
                            import pandas as pd
                            df = pd.read_excel(str(p))
                            extracted = df.head(50).to_csv(index=False)
                        except Exception:
                            extracted = ""

                    extracted = (extracted or "").strip()
                    if extracted:
                        extracted = extracted[:8000]
                        extracted_parts.append(f"=== {title or p.name} ({p.name}) ===\n{extracted}")
                        revision_context_titles.append(title or p.name)
                except Exception:
                    continue

            if extracted_parts:
                revision_context_text = "\n\n".join(extracted_parts)

        # 1. Analysis Step (Real LLM Call)
        # Update memory with current user request
        if revision_context_text:
            self.memory.append(
                {
                    "role": "system",
                    "content": "Existing deliverable content (for revision):\n" + revision_context_text,
                }
            )
        self.memory.append({"role": "user", "content": user_input})
        
        # Build context from memory for Planner
        memory_context = ""
        if len(self.memory) > 1:
            memory_context = "Previous Conversation:\n"
            for msg in self.memory[:-1]:
                memory_context += f"- {msg['role']}: {msg['content']}\n"
        
        yield {
            "type": "step_start",
            "title": "任务分析",
            "desc": "正在使用规划智能体分析用户请求...",
            "status": "running"
        }
        
        # Clarification Loop (Simulated for now, real implementation needs bidirectional)
        # We'll ask the Planner if clarification is needed.
        # But per user request, we need a Questionnaire UI if input is vague.
        # For this turn, we'll assume we proceed but we'll structure the prompt to be critical.
        
        planner_system = prompt_manager.render_prompt("deep_research/planner.system.txt", {
            "current_date": current_date
        })
        planner_user = prompt_manager.render_prompt("deep_research/planner.user.txt", {
            "user_input": user_input,
            "memory_context": memory_context,
            "current_date": current_date
        })
        
        if revision_context_text:
            planner_user += (
                "\n\nYou are revising an existing deliverable. "
                "Do not treat this as a brand-new task. "
                "Preserve the existing structure/style unless the user asks otherwise. "
                "The existing deliverable content is available in the system context."
            )
        planner_user += "\n\nIf key info is missing, return JSON {\"clarification\": {\"title\": \"...\", \"auto_decide_seconds\": 60, \"questions\": [...]}} so the UI can render a questionnaire with options. Otherwise return the plan steps."
        
        try:
            plan_json = await self._call_llm("planner", [
                {"role": "system", "content": planner_system},
                {"role": "user", "content": planner_user}
            ])
            # Simple cleanup for JSON parsing
            if "```json" in plan_json:
                plan_json = plan_json.split("```json")[1].split("```")[0].strip()
            elif "```" in plan_json:
                plan_json = plan_json.split("```")[1].split("```")[0].strip()
                
            plan_data = json.loads(plan_json)
            
            # Handle Clarification
            if isinstance(plan_data, dict) and "clarification" in plan_data:
                clarification = plan_data["clarification"]

                if isinstance(clarification, dict):
                    yield {
                        "type": "clarification",
                        "title": clarification.get("title") or "需要补充信息",
                        "auto_decide_seconds": clarification.get("auto_decide_seconds") or 60,
                        "questions": clarification.get("questions") or [],
                        "session_id": self.sandbox_session_id,
                    }
                else:
                    questions_text = clarification
                    if isinstance(questions_text, list):
                        questions_text = "\n".join([str(q) for q in questions_text if q is not None])
                    yield {
                        "type": "clarification",
                        "content": str(questions_text or ""),
                        "auto_decide_seconds": 60,
                        "session_id": self.sandbox_session_id,
                    }
                
                # In a streaming response, we yield and expect the connection to stay open or client to reconnect.
                # Ideally, we would pause here. For this implementation, we stop yielding steps.
                # The frontend will re-submit with the answer appended to the prompt.
                return 

            elif isinstance(plan_data, dict) and "steps" in plan_data:
                steps = plan_data["steps"]
            elif isinstance(plan_data, list):
                steps = plan_data
            else:
                steps = ["Search Web", "Summarize Findings", "Generate Output"]
            
            # Emit Plan Review (Todo List)
            yield { 
                "type": "plan_review", 
                "steps": ["任务分析", *steps, "报告生成"],
                "session_id": self.sandbox_session_id
            } 
            
            # Wait for approval? For now, auto-proceed, but the frontend can show the list.
            
            yield {
                "type": "log",
                "step_title": "任务分析",
                "content": f"生成计划: {', '.join(steps)}"
            }
        except Exception as e:
            yield {
                "type": "log",
                "step_title": "任务分析",
                "content": f"规划失败，使用默认计划。错误: {str(e)}"
            }
            steps = ["Search Web", "Summarize Findings", "Generate Output"]

        yield {
             "type": "step_complete",
             "title": "任务分析",
             "status": "completed"
        }

        # 2. Execution Loop (Iterate through steps)
        # We will accumulate findings in a simple string buffer for the report
        research_findings = []
        
        for step_index, step in enumerate(steps):
            yield {
                "type": "step_start",
                "title": step,
                "desc": "Executing research step...",
                "status": "running"
            }
            
            # Ask Researcher Model what to do: Search or Just Analyze?
            researcher_system = prompt_manager.render_prompt("deep_research/researcher.system.txt", {
                "current_date": current_date
            })
            researcher_user = prompt_manager.render_prompt("deep_research/researcher.user.txt", {
                "step": step,
                "current_date": current_date
            })
            
            try:
                action_response = await self._call_llm("researcher", [
                    {"role": "system", "content": researcher_system},
                    {"role": "user", "content": researcher_user}
                ])
                
                # Parse action
                action_data = {}
                try:
                    clean_json = action_response
                    if "```json" in clean_json:
                        clean_json = clean_json.split("```json")[1].split("```")[0].strip()
                    elif "```" in clean_json:
                        clean_json = clean_json.split("```")[1].split("```")[0].strip()
                    action_data = json.loads(clean_json)
                except:
                    # Fallback if LLM didn't output valid JSON
                    action_data = {"action": "analyze", "thought": action_response}

                step_content = ""
                action_type = action_data.get("action")
                
                if action_type in ["search", "academic_search", "visit"]:
                    query = action_data.get("query", step)

                    if action_type == "visit":
                        yield {
                            "type": "log",
                            "step_title": step,
                            "content": f"Deep Visiting: {query}"
                        }
                        page_text = await self.search_service.visit_page(query)
                        note_md = f"# Step: {step}\n\n## Visit\n\nURL: {query}\n\n{self._truncate_text(page_text, 20000)}\n"
                        note_path = self._write_workspace_text(f"notes/step_{step_index + 1}_visit.md", note_md)
                        step_content = f"已写入临时笔记文件: {note_path}\n\n{self._truncate_text(page_text, 3000)}"

                    elif action_type == "academic_search":
                        arxiv_ok = await self._can_access_url("https://arxiv.org/")
                        cnki_ok = await self._can_access_url("https://www.cnki.net/")
                        wos_ok = await self._can_access_url("https://www.webofscience.com/")
                        yield {
                            "type": "log",
                            "step_title": step,
                            "content": f"Searching Academic (ArXiv/CNKI/SCI) for: {query}"
                        }

                        snippets: List[str] = []
                        sources: List[Tuple[str, str]] = []

                        if arxiv_ok:
                            arxiv_has_sources = False
                            search_results = academic_search.search(query, max_results=6)
                            for r in search_results or []:
                                if not isinstance(r, dict) or "error" in r:
                                    continue
                                title = (r.get("title") or "").strip()
                                href = (r.get("href") or "").strip()
                                body = (r.get("body") or "").strip()
                                if href:
                                    arxiv_has_sources = True
                                    sources.append((title, href))
                                    yield {
                                        "type": "resource",
                                        "step_title": step,
                                        "kind": "academic_result",
                                        "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "ArXiv"},
                                    }
                                snippets.append(f"- [ArXiv] {title} ({href}): {self._truncate_text(body, 400)}")
                            if not arxiv_has_sources:
                                arxiv_payload = await self.search_service.search_structured(
                                    f"{query} site:arxiv.org",
                                    max_results=6,
                                    region="wt-wt",
                                )
                                for r in arxiv_payload.get("results") or []:
                                    title = (r.get("title") or "").strip()
                                    href = (r.get("href") or "").strip()
                                    body = (r.get("body") or "").strip()
                                    if not href:
                                        continue
                                    sources.append((title, href))
                                    yield {
                                        "type": "resource",
                                        "step_title": step,
                                        "kind": "academic_result",
                                        "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "ArXiv(Web)"},
                                    }
                                    snippets.append(f"- [ArXiv(Web)] {title} ({href}): {self._truncate_text(body, 400)}")
                        else:
                            snippets.append("- [ArXiv] (跳过：网络不可用或访问失败)")

                        if cnki_ok:
                            cnki_payload = await self.search_service.search_structured(
                                f"{query} site:cnki.net",
                                max_results=6,
                                region="zh-CN",
                            )
                            for r in cnki_payload.get("results") or []:
                                title = (r.get("title") or "").strip()
                                href = (r.get("href") or "").strip()
                                body = (r.get("body") or "").strip()
                                if href:
                                    sources.append((title, href))
                                    yield {
                                        "type": "resource",
                                        "step_title": step,
                                        "kind": "academic_result",
                                        "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "CNKI"},
                                    }
                                snippets.append(f"- [CNKI] {title} ({href}): {self._truncate_text(body, 400)}")
                        else:
                            cnki_payload = await self.search_service.search_structured(
                                f"{query} site:cnki.net",
                                max_results=6,
                                region="zh-CN",
                            )
                            for r in cnki_payload.get("results") or []:
                                title = (r.get("title") or "").strip()
                                href = (r.get("href") or "").strip()
                                body = (r.get("body") or "").strip()
                                if not href:
                                    continue
                                sources.append((title, href))
                                yield {
                                    "type": "resource",
                                    "step_title": step,
                                    "kind": "academic_result",
                                    "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "CNKI(Web)"},
                                }
                                snippets.append(f"- [CNKI(Web)] {title} ({href}): {self._truncate_text(body, 400)}")

                        if wos_ok:
                            wos_payload = await self.search_service.search_structured(
                                f"{query} site:webofscience.com",
                                max_results=6,
                                region="wt-wt",
                            )
                            for r in wos_payload.get("results") or []:
                                title = (r.get("title") or "").strip()
                                href = (r.get("href") or "").strip()
                                body = (r.get("body") or "").strip()
                                if href:
                                    sources.append((title, href))
                                    yield {
                                        "type": "resource",
                                        "step_title": step,
                                        "kind": "academic_result",
                                        "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "SCI"},
                                    }
                                snippets.append(f"- [SCI] {title} ({href}): {self._truncate_text(body, 400)}")
                        else:
                            wos_payload = await self.search_service.search_structured(
                                f"{query} site:webofscience.com",
                                max_results=6,
                                region="wt-wt",
                            )
                            for r in wos_payload.get("results") or []:
                                title = (r.get("title") or "").strip()
                                href = (r.get("href") or "").strip()
                                body = (r.get("body") or "").strip()
                                if not href:
                                    continue
                                sources.append((title, href))
                                yield {
                                    "type": "resource",
                                    "step_title": step,
                                    "kind": "academic_result",
                                    "data": {"title": title[:200], "url": href, "snippet": body[:400], "source": "SCI(Web)"},
                                }
                                snippets.append(f"- [SCI(Web)] {title} ({href}): {self._truncate_text(body, 400)}")

                        if not snippets:
                            snippets = ["No relevant results found."]

                        note_md = f"# Step: {step}\n\n## Academic Search\n\nQuery: {query}\n\n" + "\n".join(snippets) + "\n"
                        note_path = self._write_workspace_text(f"notes/step_{step_index + 1}_academic.md", note_md)
                        top_sources = "\n".join([f"- {t} ({u})" for t, u in sources[:8]]) or "- (无)"
                        step_content = f"已写入临时笔记文件: {note_path}\n\n关键来源:\n{top_sources}\n"

                    else:
                        queries = self._build_search_queries(query, user_input)
                        sources_by_url: Dict[str, Dict[str, str]] = {}
                        notes_parts: List[str] = [f"# Step: {step}\n\n## Multi-round Web Search\n"]

                        for round_index, q in enumerate(queries, start=1):
                            yield {
                                "type": "log",
                                "step_title": step,
                                "content": f"Round {round_index}/{len(queries)} 搜索: {q}"
                            }

                            yield {
                                "type": "resource",
                                "step_title": step,
                                "kind": "search",
                                "data": {"query": q, "round": round_index}
                            }

                            search_payload = await self.search_service.search_structured(q, max_results=6)
                            search_text = search_payload.get("formatted") or ""
                            results = search_payload.get("results") or []

                            notes_parts.append(f"### Round {round_index}: {q}\n\n{search_text}\n")

                            for r in results:
                                href = (r.get("href") or "").strip()
                                if not href:
                                    continue
                                if href not in sources_by_url:
                                    sources_by_url[href] = {
                                        "title": (r.get("title") or "").strip(),
                                        "url": href,
                                        "snippet": (r.get("body") or "").strip(),
                                    }
                                yield {
                                    "type": "resource",
                                    "step_title": step,
                                    "kind": "search_result",
                                    "data": {
                                        "title": (r.get("title", "") or "")[:200],
                                        "url": href,
                                        "snippet": (r.get("body", "") or "")[:400],
                                        "round": round_index
                                    }
                                }

                            visited_pages: List[str] = []
                            for r in results[:2]:
                                href = (r.get("href") or "").strip()
                                if not href:
                                    continue
                                yield {
                                    "type": "resource",
                                    "step_title": step,
                                    "kind": "visit",
                                    "data": {
                                        "title": (r.get("title", "") or "")[:200],
                                        "url": href,
                                        "snippet": "访问网页并提取正文...",
                                        "round": round_index
                                    }
                                }
                                page_text = await self.search_service.visit_page(href)
                                visited_pages.append(f"#### Visited: {href}\n\n{self._truncate_text(page_text, 12000)}\n")

                            if visited_pages:
                                notes_parts.append("### Visited Pages\n\n" + "\n\n".join(visited_pages))

                            image_urls = search_payload.get("images") or re.findall(r"\[IMAGE: (.*?)\]", search_text)
                            if image_urls:
                                yield {
                                    "type": "log",
                                    "step_title": step,
                                    "content": f"Round {round_index}: 发现 {len(image_urls)} 张图片，分析前 2 张..."
                                }
                                visual_parts: List[str] = []
                                for img_url in image_urls[:2]:
                                    try:
                                        description = await vision_service.describe_image(
                                            img_url,
                                            prompt="Describe this image relevant to the research topic."
                                        )
                                        visual_parts.append(f"- {img_url}\n\n{self._truncate_text(description, 1200)}\n")
                                    except Exception:
                                        continue
                                if visual_parts:
                                    notes_parts.append("### Visual Analysis\n\n" + "\n".join(visual_parts))

                        sources_lines = "\n".join(
                            [
                                f"- {self._truncate_text(v.get('title') or v.get('url') or '', 120)} ({v.get('url')})"
                                for v in list(sources_by_url.values())[:12]
                            ]
                        ) or "- (无)"

                        notes_parts.append("## Key Sources\n\n" + sources_lines + "\n")
                        note_md = "\n\n".join(notes_parts).strip() + "\n"
                        note_path = self._write_workspace_text(f"notes/step_{step_index + 1}_web.md", note_md)

                        step_content = f"已写入临时笔记文件: {note_path}\n\n关键来源:\n{sources_lines}\n"

                    yield {
                        "type": "log",
                        "step_title": step,
                        "content": "Search and analysis completed."
                    }
                    
                elif action_type == "python_sandbox":
                    code = action_data.get("query", "") # prompt uses 'query' field for code
                    yield {
                        "type": "log",
                        "step_title": step,
                        "content": f"Executing in Sandbox: {code[:50]}..."
                    }
                    result = self.execute_code(code)
                    step_content = f"Sandbox Output:\n{result.get('output', '')}\nError: {result.get('error', '')}"
                
                elif action_type == "python_chart":
                    chart_code = action_data.get("code") or action_data.get("query", "")
                    yield {
                        "type": "log",
                        "step_title": step,
                        "content": f"Generating chart with code: {chart_code[:100]}..."
                    }
                    # We inject code to print the data as CSV so we can capture it
                    # We also try to print the dataframe directly if it exists in local vars, but explicit print is safer
                    wrapped_code = f"{chart_code}\nimport pandas as pd\nif 'data' in locals():\n    if isinstance(data, dict): print('---CHART_DATA_START---'); print(pd.DataFrame(data).to_csv(index=False)); print('---CHART_DATA_END---')\n    elif isinstance(data, pd.DataFrame): print('---CHART_DATA_START---'); print(data.to_csv(index=False)); print('---CHART_DATA_END---')"
                    
                    code_result = self.execute_code(wrapped_code)
                    
                    if code_result.get("success") and code_result.get("output"):
                        output = code_result['output']
                        # Extract CSV part if marked
                        if "---CHART_DATA_START---" in output:
                            csv_content = output.split("---CHART_DATA_START---")[1].split("---CHART_DATA_END---")[0].strip()
                            step_content += f"\n\n[CHART_DATA_CSV]\n{csv_content}\n[/CHART_DATA_CSV]"
                            yield {
                                "type": "log",
                                "step_title": step,
                                "content": "Chart data generated successfully."
                            }
                        else:
                            # Fallback if they didn't use 'data' variable or just printed something else
                            step_content += f"\n\nCode Output:\n{output}"
                    else:
                        step_content += f"\n\nChart generation failed: {code_result.get('error', 'Unknown error')}"

                else:
                    thought = action_data.get("thought", "Analyzing data...")
                    yield {
                        "type": "log", 
                        "step_title": step,
                        "content": thought
                    }
                    step_content = f"Analysis: {thought}"

                # Store findings
                research_findings.append(f"## Step: {step}\n{step_content}\n")

            except Exception as e:
                yield {
                    "type": "log", 
                    "step_title": step,
                    "content": f"Error executing step: {str(e)}"
                }
                
            yield {
                 "type": "step_complete",
                 "title": step,
                 "status": "completed"
            }

        # 3. Final Artifact Generation (Writer Agent)
        yield {
            "type": "step_start",
            "title": "报告生成",
            "desc": "Compiling final report...",
            "status": "running"
        }
        
        # Generate Report Content using Writer LLM
        all_findings_text = "\n".join(research_findings)
        writer_system = prompt_manager.render_prompt("deep_research/writer.system.txt", {
            "current_date": current_date
        })
        writer_user = prompt_manager.render_prompt("deep_research/writer.user.txt", {
            "user_input": user_input,
            "research_findings": all_findings_text,
            "current_date": current_date
        })
        if revision_context_text:
            titles_text = ", ".join([t for t in revision_context_titles if t]) or "existing file"
            writer_user += (
                "\n\nIMPORTANT: This is a revision request. "
                f"You must modify the existing deliverable(s): {titles_text}. "
                "Do not restate the user request text. "
                "Keep structure and styling consistent unless explicitly requested to change."
            )
        
        # Inject instruction for images
        # writer_user += "\n\nIMPORTANT: If there are image URLs or Visual Analysis in the findings, please include the most relevant images in the report using HTML <img> tags (e.g., <img src='url'>)."
        # writer_user += "\n\nIMPORTANT: If PPT is requested, the PPT 'content' must be HTML with multiple <div class=\"ppt-slide\"> blocks. If a bar chart is needed, include at least one <table class=\"chart\">."
        # writer_user += "\n\nIMPORTANT: Return a JSON object with 'files': [{'format': 'docx', 'content': '...'}, {'format': 'ppt', 'content': '...'}] if multiple formats are requested."

        report_failed = False
        try:
            writer_response = await self._call_llm("writer", [
                {"role": "system", "content": writer_system},
                {"role": "user", "content": writer_user}
            ])

            output_data_list = self._parse_writer_output_files(writer_response, user_input)

            requested_formats = self._detect_requested_formats(user_input)
            produced_formats = {self._normalize_output_format(item.get("format")) for item in output_data_list}
            missing_formats = [f for f in requested_formats if f not in produced_formats]

            if missing_formats:
                fix_user = (
                    writer_user
                    + "\n\nYou omitted required deliverables. "
                    + f"Return JSON with files for formats: {missing_formats}. "
                    + "Keep content consistent with prior answer. "
                    + "For PPT, content must be HTML with multiple <div class=\"ppt-slide\"> blocks."
                )
                writer_response_2 = await self._call_llm("writer", [
                    {"role": "system", "content": writer_system},
                    {"role": "user", "content": fix_user}
                ])
                output_data_list = self._parse_writer_output_files(writer_response_2, user_input)

            produced_formats = {self._normalize_output_format(item.get("format")) for item in output_data_list}
            missing_formats = [f for f in requested_formats if f not in produced_formats]
            if missing_formats:
                filename_base = f"Report_{int(datetime.datetime.now().timestamp())}"
                for f in missing_formats:
                    if f == "ppt":
                        output_data_list.append({
                            "format": "ppt",
                            "content": self._build_fallback_ppt_html(user_input, all_findings_text),
                            "filename": f"{filename_base}.pptx",
                        })
                    elif f == "docx":
                        output_data_list.append({
                            "format": "docx",
                            "content": self._ensure_html(all_findings_text),
                            "filename": f"{filename_base}.docx",
                        })
                    elif f == "pdf":
                        output_data_list.append({
                            "format": "pdf",
                            "content": self._ensure_html(all_findings_text),
                            "filename": f"{filename_base}.pdf",
                        })
                    elif f == "excel":
                        output_data_list.append({
                            "format": "excel",
                            "content": [{"说明": "未返回可用的结构化表格数据"}],
                            "filename": f"{filename_base}.xlsx",
                        })

            for output_data in output_data_list:
                output_format = self._normalize_output_format(output_data.get("format", "docx"))
                content = output_data.get("content")
                filename = output_data.get("filename", f"Report_{int(datetime.datetime.now().timestamp())}")
                
                output_dir = "app/static/reports" # Fixed path relative to backend root
                os.makedirs(output_dir, exist_ok=True)
                
                artifact_path = ""
                artifact_type = "DOCX"
                
                if output_format == "excel" and isinstance(content, list):
                    artifact_type = "XLSX"
                    if not filename.endswith(".xlsx"): filename += ".xlsx"
                    excel_generator.output_dir = output_dir
                    try:
                        artifact_path = excel_generator.generate_excel(content, filename)
                    except Exception:
                        skywork_query = f"{user_input}\n\n{all_findings_text}"
                        skywork_url = await self._skywork_generate_file_url("sheet", skywork_query)
                        if not skywork_url:
                            raise
                        artifact_path = await self._download_to_path(skywork_url, os.path.join(output_dir, filename))
                    
                elif output_format == "ppt":
                    artifact_type = "PPTX"
                    if not filename.endswith(".pptx"): filename += ".pptx"
                    ppt_html = str(content) if isinstance(content, str) else str(content)
                    yield {
                        "type": "artifact_preview",
                        "data": {
                            "title": filename,
                            "format": "ppt",
                            "html": self._wrap_preview_html("ppt", ppt_html),
                        },
                    }
                    try:
                        artifact_path = self.ppt_generator.generate_ppt(ppt_html, filename)
                    except Exception:
                        skywork_query = f"{user_input}\n\n{all_findings_text}"
                        skywork_url = await self._skywork_generate_file_url("ppt", skywork_query)
                        if not skywork_url:
                            raise
                        artifact_path = await self._download_to_path(skywork_url, os.path.join(output_dir, filename))

                elif output_format == "pdf":
                    artifact_type = "PDF"
                    if not filename.endswith(".pdf"): filename += ".pdf"

                    # Convert content to string/HTML if needed
                    if not isinstance(content, str):
                        if isinstance(content, list):
                            html_body = ""
                            for section in content:
                                title = section.get("title", "")
                                body = section.get("content", "")
                                if isinstance(body, list):
                                    body = "<ul>" + "".join([f"<li>{item}</li>" for item in body]) + "</ul>"
                                html_body += f"<h2>{title}</h2><div>{body}</div>"
                            content = html_body
                        else:
                            content = str(content)
                    else:
                        content = self._ensure_html(content)

                    yield {
                        "type": "artifact_preview",
                        "data": {
                            "title": filename,
                            "format": "pdf",
                            "html": self._wrap_preview_html("doc", content),
                        },
                    }
                    try:
                        artifact_path = self.pdf_generator.generate_pdf(content, filename)
                        if isinstance(artifact_path, str) and artifact_path.startswith("Error:"):
                            raise RuntimeError(artifact_path)
                    except Exception:
                        skywork_query = f"{user_input}\n\n{all_findings_text}"
                        skywork_url = await self._skywork_generate_file_url("pdf", skywork_query)
                        if not skywork_url:
                            raise
                        artifact_path = await self._download_to_path(skywork_url, os.path.join(output_dir, filename))

                else: # Default to DOCX
                    artifact_type = "DOCX"
                    if not filename.endswith(".docx"): filename += ".docx"
                    
                    # Ensure content is string
                    if not isinstance(content, str): content = str(content)
                    content = self._ensure_html(content)
                    yield {
                        "type": "artifact_preview",
                        "data": {
                            "title": filename,
                            "format": "docx",
                            "html": self._wrap_preview_html("doc", content),
                        },
                    }
                    try:
                        artifact_path = self.word_generator.generate_docx(content, filename)
                    except Exception:
                        skywork_query = f"{user_input}\n\n{all_findings_text}"
                        skywork_url = await self._skywork_generate_file_url("doc", skywork_query)
                        if not skywork_url:
                            raise
                        artifact_path = await self._download_to_path(skywork_url, os.path.join(output_dir, filename))
                
                # Get file size
                if not artifact_path or not os.path.exists(artifact_path):
                    raise RuntimeError(f"Artifact generation failed: {artifact_path}")
                file_size = os.path.getsize(artifact_path)
                size_str = f"{file_size / 1024:.1f} KB"
                
                yield {
                    "type": "artifact",
                    "data": {
                        "title": filename,
                        "type": artifact_type,
                        "size": size_str,
                        "path": artifact_path
                    }
                }
            
        except Exception as e:
            report_failed = True
            yield {
                "type": "log",
                "step_title": "报告生成",
                "content": f"生成报告失败: {str(e)}"
            }
        
        yield {
            "type": "step_complete",
            "title": "报告生成",
            "status": "completed"
        }

        yield {
            "type": "final_result",
            "content": "研究任务已完成。" if report_failed else "研究任务成功完成。"
        }

    def execute_code(self, code: str) -> Dict[str, Any]:
        """Runs code in the isolated sandbox."""
        return sandbox_service.execute_code(self.sandbox_session_id, code)

    def generate_ppt(self, html: str, filename: str) -> str:
        return self.ppt_generator.generate_ppt(html, filename)
        
    def generate_word(self, html: str, filename: str) -> str:
        return self.word_generator.generate_docx(html, filename)

    def cleanup(self):
        """Destroys the sandbox session."""
        sandbox_service.delete_session(self.sandbox_session_id)
