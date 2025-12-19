"""
File Ingestion Service
Handles parsing of various file formats (PDF, DOCX, Images) and integrates with Vision Agents.
"""

import os
from typing import Dict, Any, List, Optional
from pathlib import Path
import io

# Import our new office tools
from app.tools.office_suite import OfficeProcessor
from docx import Document
import pdfplumber
from pptx import Presentation
from PIL import Image

class FileIngestionService:
    def __init__(self, upload_dir: str = "uploads"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.processor = OfficeProcessor(output_dir=str(self.upload_dir))

    async def parse_file(self, file_path: str, vision_agent_callback=None) -> Dict[str, Any]:
        """
        Parses a file and returns structured content.
        Supports PDF, DOCX, PPTX image extraction.
        
        Args:
            file_path: Absolute path to the file.
            vision_agent_callback: Async function to call for image description. 
                                   Signature: async (image_path) -> str
        """
        path = Path(file_path)
        ext = path.suffix.lower()
        
        content = ""
        metadata = {"filename": path.name, "type": ext, "images_found": 0}

        try:
            if ext == ".pdf":
                # 1. Extract Text
                text = ""
                with pdfplumber.open(path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text += f"\n--- Page {i+1} ---\n"
                        text += page.extract_text() or ""
                        
                        # 2. Extract Images (Experimental)
                        if vision_agent_callback:
                            # Note: pdfplumber extracts image objects, converting to file for Vision API
                            # This is complex as it returns raw bytes or objects.
                            # For simplicity in this MVP, we skip complex PDF image extraction 
                            # unless we want to use `page.images` and crop.
                            pass 
                            
                content += f"--- PDF CONTENT ---\n{text}\n"
                
            elif ext == ".docx":
                doc = Document(str(path))
                full_text = []
                
                # Text
                for para in doc.paragraphs:
                    full_text.append(para.text)
                
                # Images (rIds)
                # DOCX images are stored in relationships. 
                # python-docx doesn't easily expose image data in order with text without complex XML parsing.
                # We will extract all images found in the package.
                if vision_agent_callback:
                    for rel in doc.part.rels.values():
                        if "image" in rel.target_ref:
                            # This is a bit low-level, skipping for stability in MVP
                            pass

                content += f"--- DOCX CONTENT ---\n" + "\n".join(full_text)
                
            elif ext == ".pptx":
                prs = Presentation(str(path))
                text_content = []
                
                for i, slide in enumerate(prs.slides):
                    text_content.append(f"\n--- Slide {i+1} ---")
                    
                    # Text and Images
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                        
                        if shape.shape_type == 13: # PICTURE
                            if vision_agent_callback:
                                # Save image to temp
                                image = shape.image
                                image_bytes = image.blob
                                img_ext = image.ext
                                img_filename = f"{path.stem}_slide{i+1}_{shape.shape_id}.{img_ext}"
                                img_path = self.upload_dir / "extracted_images" / img_filename
                                img_path.parent.mkdir(exist_ok=True)
                                
                                with open(img_path, "wb") as f:
                                    f.write(image_bytes)
                                    
                                # Call Vision Agent
                                desc = await vision_agent_callback(str(img_path))
                                text_content.append(f"[Image on Slide {i+1}: {desc}]")
                                metadata["images_found"] += 1

                content += f"--- PPTX CONTENT ---\n" + "\n".join(text_content)

            elif ext in [".txt", ".md", ".py", ".json", ".csv", ".html", ".js", ".css", ".xml", ".yaml", ".yml"]:
                with open(path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                    
            elif ext in [".odt", ".ods", ".odp"]:
                 # Placeholder for OpenOffice formats (would require `odfpy` or similar)
                 # For now, we treat them as unsupported or try text extraction if possible
                 return {"error": f"OpenOffice format {ext} support is coming soon."}

            elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
                if vision_agent_callback:
                    description = await vision_agent_callback(str(path))
                    content = f"--- IMAGE DESCRIPTION ---\n{description}"
                else:
                    content = "[Image file uploaded, but no Vision Agent available to describe it.]"
            
            else:
                return {"error": f"Unsupported file type: {ext}"}
                
        except Exception as e:
            return {"error": f"Failed to parse file: {str(e)}"}

        return {
            "content": content,
            "metadata": metadata
        }

file_ingestion_service = FileIngestionService()
